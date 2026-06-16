"""
drive_content.py — Lectura real del contenido de los archivos de Drive.

Estrategia (texto primero, PDF nativo solo si no hay otra opción):
  1. Si el archivo tiene texto extraíble (PDF con texto, DOCX, XLSX, Google Docs,
     Google Sheets, texto plano, ZIP-WhatsApp) → se extrae el texto y se pasa como
     bloque de texto. Mucho más barato en tokens = caben más archivos por batch.
  2. Si la extracción de texto falla, retorna vacío o lanza excepción → se pasa el
     archivo NATIVO como document/image block para que Claude lo lea con visión
     (PDFs escaneados, imágenes, diagramas, clippings de prensa).
  3. Las imágenes siempre se pasan como image block (no tienen texto extraíble).

Esto permite leer entre 3-5× más archivos por llamada comparado con pasar todo
como PDFs nativos en base64.
"""
from __future__ import annotations

import base64
import csv
import io
import logging
import zipfile

logger = logging.getLogger(__name__)

# Mínimo de texto "útil" para considerar que la extracción tuvo éxito.
# Menos que esto = probablemente un PDF escaneado → pasar como nativo.
MIN_TEXT_CHARS = 150

MAX_BINARY_BYTES = 20 * 1024 * 1024      # PDFs / imágenes individuales (fallback nativo)
MAX_TEXT_CHARS = 60_000                   # caracteres por archivo extraído a texto
MAX_TOTAL_BINARY_BYTES = 28 * 1024 * 1024  # presupuesto binario acumulado por cuenta
# Tope de DESCARGA para PDFs grandes: aunque no quepan como bloque nativo,
# vale la pena bajarlos para extraerles el texto (gratis en tokens) o
# rasterizar sus primeras páginas si son escaneados.
MAX_PDF_DOWNLOAD_BYTES = 100 * 1024 * 1024
# Rasterización de PDFs escaneados grandes: páginas máximas y resolución.
RASTER_MAX_PAGES = 6
RASTER_DPI = 110

# ── MIME types ────────────────────────────────────────────────────────────────
PDF_MIME = "application/pdf"
GOOGLE_DOC = "application/vnd.google-apps.document"
GOOGLE_SLIDES = "application/vnd.google-apps.presentation"
GOOGLE_SHEET = "application/vnd.google-apps.spreadsheet"
GOOGLE_DRAWING = "application/vnd.google-apps.drawing"
GOOGLE_SHORTCUT = "application/vnd.google-apps.shortcut"
FOLDER_MIME = "application/vnd.google-apps.folder"

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

IMAGE_MIMES = {"image/jpeg", "image/png", "image/gif", "image/webp"}

TEXT_MIMES = {
    "text/plain", "text/csv", "text/markdown", "text/tab-separated-values",
    "application/json", "text/html", "application/xml", "text/xml",
    "application/rtf", "text/rtf", "text/vcard", "text/x-vcard",
}

# Correos y páginas web guardadas (MIME multipart) — se parsean con email stdlib
EMAIL_LIKE_MIMES = {"message/rfc822", "multipart/related"}

# Subcarpetas del playbook ordenadas por prioridad de lectura, y cuántos
# archivos (más recientes) leer de cada una.
# 01 va PRIMERO: el contrato/plan de trabajo es la base de todos los compromisos.
SUBFOLDER_READ_BUDGET = {
    "01": 8,    # Contrato / OC / Plan de trabajo — PRIMERO para extraer compromisos
    "02": 14,   # Entregables — evidencia del trabajo realizado
    "03": 10,   # Reportes de avance
    "05": 8,    # Transcripciones / minutas de llamadas
    "04": 8,    # Conversaciones WhatsApp (chat + fotos)
    "06": 5,    # Agenda / briefs / calendarios
}
# Tope global de archivos a leer por cuenta.
# Para el baseline inicial subimos a 50 para capturar más archivos por cuenta.
MAX_FILES_TO_READ = 50


# ─────────────────────────────────────────────────────────────────────────────
# Selección de archivos a leer
# ─────────────────────────────────────────────────────────────────────────────

def select_files_for_analysis(drive_files: list[dict]) -> list[dict]:
    """
    Prioriza qué archivos leer: los más recientes de cada subcarpeta del playbook,
    según SUBFOLDER_READ_BUDGET, respetando el tope global MAX_FILES_TO_READ.
    """
    if not drive_files:
        return []

    by_slot: dict[str, list[dict]] = {}
    for f in drive_files:
        slot = (f.get("subfolder") or "??")
        by_slot.setdefault(slot, []).append(f)

    for slot in by_slot:
        by_slot[slot].sort(key=lambda x: x.get("modifiedTime") or "", reverse=True)

    selected: list[dict] = []
    # Primero, la cuota por subcarpeta priorizada
    for slot, budget in SUBFOLDER_READ_BUDGET.items():
        selected.extend(by_slot.get(slot, [])[:budget])

    # Luego rellenar con cualquier otro archivo reciente que no esté ya incluido
    if len(selected) < MAX_FILES_TO_READ:
        already = {f.get("id") for f in selected}
        leftovers = [f for f in drive_files if f.get("id") not in already]
        leftovers.sort(key=lambda x: x.get("modifiedTime") or "", reverse=True)
        selected.extend(leftovers[: MAX_FILES_TO_READ - len(selected)])

    return selected[:MAX_FILES_TO_READ]


# ─────────────────────────────────────────────────────────────────────────────
# Construcción de content blocks
# ─────────────────────────────────────────────────────────────────────────────

def build_content_blocks(
    files_api,
    files: list[dict],
    *,
    wa_context: dict | None = None,
) -> tuple[list[dict], list[str], str | None]:
    """
    Descarga cada archivo y lo empaca como una "unidad" de contenido.

    Cada unidad agrupa los content blocks de UN archivo (etiqueta + contenido)
    junto con su descriptor, para que el analizador pueda descartar archivos
    completos al ajustar el payload al límite de tokens (ITPM).

    Args:
        wa_context: cuando se provee, los ZIPs de la subcarpeta 04 (WhatsApp)
            se extraen de forma INCREMENTAL (solo mensajes nuevos desde el
            watermark). Formato: {"watermark_iso": str|None, "rolling_summary": str|None}

    Returns:
        units:          [{"blocks": [...], "read_file": {...}}], en orden de prioridad.
        notes:          avisos de archivos omitidos.
        latest_wa_ts:   ISO del mensaje WA más reciente visto en esta corrida
                        (None si no se procesó ningún ZIP incremental).
    """
    units: list[dict] = []
    notes: list[str] = []
    total_binary = 0
    latest_wa_ts: str | None = None

    for f in files:
        name = f.get("name") or "(sin nombre)"
        mime = f.get("mimeType") or ""
        fid = f.get("id")
        slot = f.get("subfolder") or "??"
        label = f"[{slot}] {name}"
        if not fid:
            continue

        # ── WhatsApp ZIP incremental (subcarpeta 04) ──────────────────────────
        is_zip = mime == "application/zip" or name.lower().endswith(".zip")
        if is_zip and slot == "04" and wa_context is not None:
            try:
                import wa_parser
                data = _download(files_api.get_media(fileId=fid))
                text, wa_ts = wa_parser.extract_incremental(
                    data,
                    wa_context.get("watermark_iso"),
                    wa_context.get("rolling_summary"),
                )
                if wa_ts:
                    latest_wa_ts = wa_ts
                block = _text_block(text)
                kind = "WhatsApp/incremental"
                logger.info(
                    "    WA incremental: %s | watermark=%s → latest=%s",
                    name,
                    wa_context.get("watermark_iso") or "ninguno",
                    wa_ts or "?",
                )
            except Exception as e:  # noqa: BLE001
                logger.warning("    WA incremental falló para %s: %s — usando extracción completa", name, e)
                try:
                    data = _download(files_api.get_media(fileId=fid))
                    block = _text_block(_extract_zip_text_raw(data))
                    kind = "WhatsApp/ZIP→texto (fallback)"
                except Exception as e2:
                    logger.warning("    No se pudo leer %s: %s", name, e2)
                    notes.append(f"{label}: error de lectura ({type(e2).__name__})")
                    continue

            if block is None:
                notes.append(f"{label}: sin contenido extraíble")
                continue

            units.append({
                "blocks": [
                    {"type": "text", "text": f"\n----- ARCHIVO: {label} ({kind}) -----"},
                    block,
                ],
                "read_file": {
                    "id": fid,
                    "title": name,
                    "subfolder": slot,
                    "subfolderName": f.get("subfolderName"),
                    "modifiedTime": f.get("modifiedTime"),
                    "kind": kind,
                },
            })
            continue

        # ── Resto de archivos: flujo estándar ─────────────────────────────────
        try:
            block, kind, room = _file_to_block(files_api, f, total_binary)
        except Exception as e:  # noqa: BLE001 — un archivo no debe tumbar la cuenta
            logger.warning("    No se pudo leer %s (%s): %s", name, mime, e)
            notes.append(f"{label}: error de lectura ({type(e).__name__})")
            continue

        if block is None:
            notes.append(f"{label}: {kind}")
            continue

        total_binary += room
        # block puede ser un dict o una lista de dicts (p.ej. PDF rasterizado
        # a varias imágenes de página)
        content_blocks = block if isinstance(block, list) else [block]
        units.append({
            "blocks": [
                {"type": "text", "text": f"\n----- ARCHIVO: {label} ({kind}) -----"},
                *content_blocks,
            ],
            "read_file": {
                "id": fid,
                "title": name,
                "subfolder": slot,
                "subfolderName": f.get("subfolderName"),
                "modifiedTime": f.get("modifiedTime"),
                "kind": kind,
            },
        })

    return units, notes, latest_wa_ts


def _file_to_block(files_api, f: dict, total_binary: int) -> tuple[dict | list[dict] | None, str, int]:
    """
    Convierte un archivo en un content block siguiendo la estrategia texto-primero:

    1. Intenta extraer texto (pypdf/pymupdf para PDFs, python-docx para Word, etc.)
    2. Si el texto es suficientemente rico (≥ MIN_TEXT_CHARS) → bloque de texto (barato)
    3. Si la extracción falla o retorna poco texto → PDF/imagen nativo (costoso pero completo)

    Returns (block | lista de blocks | None, kind_or_reason, binary_bytes_used).
    """
    mime = f.get("mimeType") or ""
    name = (f.get("name") or "").lower()
    fid = f["id"]
    size = _to_int(f.get("size"))

    # ── Atajos de Drive → resolver el archivo destino y leerlo ────────────────
    if mime == GOOGLE_SHORTCUT:
        meta = files_api.get(
            fileId=fid, fields="shortcutDetails", supportsAllDrives=True
        ).execute(num_retries=3)
        details = (meta or {}).get("shortcutDetails") or {}
        target_id = details.get("targetId")
        target_mime = details.get("targetMimeType") or ""
        if not target_id:
            return None, "atajo sin destino", 0
        if target_mime == FOLDER_MIME:
            return None, "atajo a carpeta (no legible)", 0
        resolved = {**f, "id": target_id, "mimeType": target_mime, "size": None}
        block, kind, room = _file_to_block(files_api, resolved, total_binary)
        return block, f"atajo→{kind}", room

    # ── Google Docs → texto plano (mucho más barato que PDF export) ───────────
    if mime == GOOGLE_DOC:
        try:
            data = _download(files_api.export_media(fileId=fid, mimeType="text/plain"))
            text = data.decode("utf-8", errors="replace").strip()
            if len(text) >= MIN_TEXT_CHARS:
                return _text_block(text), "Google Doc→texto", 0
        except Exception:
            pass
        # Fallback: exportar como PDF nativo
        data = _download(files_api.export_media(fileId=fid, mimeType=PDF_MIME))
        if total_binary + len(data) > MAX_TOTAL_BINARY_BYTES:
            return None, "Google Doc PDF excede presupuesto", 0
        return _pdf_block(data, f.get("name")), "Google Doc→PDF (escaneado)", len(data)

    # ── Google Slides → texto plano ────────────────────────────────────────────
    if mime == GOOGLE_SLIDES:
        try:
            data = _download(files_api.export_media(fileId=fid, mimeType="text/plain"))
            text = data.decode("utf-8", errors="replace").strip()
            if len(text) >= MIN_TEXT_CHARS:
                return _text_block(text), "Google Slides→texto", 0
        except Exception:
            pass
        data = _download(files_api.export_media(fileId=fid, mimeType=PDF_MIME))
        if total_binary + len(data) > MAX_TOTAL_BINARY_BYTES:
            return None, "Slides PDF excede presupuesto", 0
        return _pdf_block(data, f.get("name")), "Google Slides→PDF", len(data)

    # ── Google Sheets → CSV ────────────────────────────────────────────────────
    if mime == GOOGLE_SHEET:
        data = _download(files_api.export_media(fileId=fid, mimeType="text/csv"))
        text = data.decode("utf-8", errors="replace")
        return _text_block(text), "Google Sheet→CSV", 0

    # ── Google Drawing → PNG (siempre imagen) ─────────────────────────────────
    if mime == GOOGLE_DRAWING:
        data = _download(files_api.export_media(fileId=fid, mimeType="image/png"))
        if total_binary + len(data) > MAX_TOTAL_BINARY_BYTES:
            return None, "Drawing excede presupuesto", 0
        return _image_block(data, "image/png"), "Google Drawing→imagen", len(data)

    # ── PDF: intentar extracción de texto primero ──────────────────────────────
    if mime == PDF_MIME or name.endswith(".pdf"):
        # Antes se rechazaban los PDFs >20MB sin descargarlos. Ahora los bajamos
        # (hasta 100MB): la extracción de texto es gratis en tokens, y si son
        # escaneados rasterizamos sus primeras páginas como imágenes.
        if size and size > MAX_PDF_DOWNLOAD_BYTES:
            return None, f"PDF demasiado grande ({_mb(size)})", 0
        data = _download(files_api.get_media(fileId=fid))
        # PDFs gigantes (clippings de prensa de 50MB): cap de texto más corto
        # para que no acaparen el presupuesto de tokens de la cuenta.
        text_cap = 20_000 if len(data) > MAX_BINARY_BYTES else MAX_TEXT_CHARS
        # Intento 1: pymupdf (fitz) — el más robusto
        text = _extract_pdf_text_fitz(data)
        if len(text) >= MIN_TEXT_CHARS:
            logger.debug("      PDF→texto (fitz) %s [%d chars]", f.get("name", ""), len(text))
            return _text_block(text[:text_cap]), "PDF→texto", 0
        # Intento 2: pdfminer
        text = _extract_pdf_text_pdfminer(data)
        if len(text) >= MIN_TEXT_CHARS:
            logger.debug("      PDF→texto (pdfminer) %s [%d chars]", f.get("name", ""), len(text))
            return _text_block(text[:text_cap]), "PDF→texto", 0
        # Fallback 1: PDF escaneado chico → pasar nativo para visión
        if len(data) <= MAX_BINARY_BYTES and total_binary + len(data) <= MAX_TOTAL_BINARY_BYTES:
            logger.debug("      PDF→nativo (escaneado) %s", f.get("name", ""))
            return _pdf_block(data, f.get("name")), "PDF escaneado→nativo", len(data)
        # Fallback 2: PDF escaneado grande → rasterizar primeras páginas a JPEG
        pages = _rasterize_pdf_pages(data)
        raster_bytes = sum(len(p) for p in pages)
        if pages and total_binary + raster_bytes <= MAX_TOTAL_BINARY_BYTES:
            logger.debug(
                "      PDF→imágenes (%d págs, %s) %s",
                len(pages), _mb(raster_bytes), f.get("name", ""),
            )
            blocks = [_image_block(p, "image/jpeg") for p in pages]
            return blocks, f"PDF escaneado→imágenes ({len(pages)} págs)", raster_bytes
        return None, f"PDF escaneado excede presupuesto ({_mb(len(data))})", 0

    # ── Imágenes → siempre nativo (no hay texto que extraer) ──────────────────
    if mime in IMAGE_MIMES:
        data = _download(files_api.get_media(fileId=fid))
        if len(data) > MAX_BINARY_BYTES or total_binary + len(data) > MAX_TOTAL_BINARY_BYTES:
            return None, "imagen excede presupuesto", 0
        return _image_block(data, mime), "imagen", len(data)

    # ── Office Word → texto ────────────────────────────────────────────────────
    if mime == DOCX_MIME or name.endswith(".docx"):
        data = _download(files_api.get_media(fileId=fid))
        return _text_block(_extract_docx(data)), "Word→texto", 0

    if mime == PPTX_MIME or name.endswith(".pptx"):
        data = _download(files_api.get_media(fileId=fid))
        return _text_block(_extract_pptx(data)), "PowerPoint→texto", 0

    if mime == XLSX_MIME or name.endswith(".xlsx"):
        data = _download(files_api.get_media(fileId=fid))
        return _text_block(_extract_xlsx(data)), "Excel→texto", 0

    # ── WhatsApp export (.zip) ─────────────────────────────────────────────────
    if mime == "application/zip" or name.endswith(".zip"):
        data = _download(files_api.get_media(fileId=fid))
        return _text_block(_extract_zip_text(data)), "WhatsApp/ZIP→texto", 0

    # ── Correos (.eml) y páginas guardadas (.mhtml) → texto ───────────────────
    if mime in EMAIL_LIKE_MIMES or name.endswith((".eml", ".mhtml", ".mht")):
        data = _download(files_api.get_media(fileId=fid))
        return _text_block(_extract_email_like(data)), "correo/mhtml→texto", 0

    # ── Texto plano y derivados (incluye tarjetas de contacto .vcf) ───────────
    if mime in TEXT_MIMES or name.endswith((".txt", ".csv", ".md", ".json", ".html", ".xml", ".tsv", ".rtf", ".vtt", ".srt", ".vcf")):
        data = _download(files_api.get_media(fileId=fid))
        return _text_block(data.decode("utf-8", errors="replace")), "texto", 0

    # ── Formatos no soportados ─────────────────────────────────────────────────
    return None, f"formato no soportado ({mime or 'desconocido'})", 0


# ─────────────────────────────────────────────────────────────────────────────
# Builders de blocks
# ─────────────────────────────────────────────────────────────────────────────

def _pdf_block(data: bytes, title: str | None) -> dict:
    return {
        "type": "document",
        "source": {
            "type": "base64",
            "media_type": PDF_MIME,
            "data": base64.standard_b64encode(data).decode("ascii"),
        },
        "title": (title or "documento")[:200],
    }


def _image_block(data: bytes, media_type: str) -> dict:
    return {
        "type": "image",
        "source": {
            "type": "base64",
            "media_type": media_type,
            "data": base64.standard_b64encode(data).decode("ascii"),
        },
    }


def _text_block(text: str) -> dict | None:
    text = (text or "").strip()
    if not text:
        return None
    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS] + "\n…[contenido truncado por longitud]…"
    return {"type": "text", "text": text}


# ─────────────────────────────────────────────────────────────────────────────
# Descarga
# ─────────────────────────────────────────────────────────────────────────────

def _download(request) -> bytes:
    """Descarga (o exporta) un archivo de Drive a bytes en memoria."""
    import time

    from googleapiclient.http import MediaIoBaseDownload

    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, request, chunksize=4 * 1024 * 1024)
    done = False
    while not done:
        # Reintento propio: errores de socket en Windows (WinError 10054)
        # no siempre los cubre num_retries de googleapiclient.
        for attempt in range(1, 5):
            try:
                _, done = downloader.next_chunk(num_retries=3)
                break
            except Exception:
                if attempt == 4:
                    raise
                time.sleep(2 * attempt)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# Extractores de texto para PDF (texto primero, nativo si falla)
# ─────────────────────────────────────────────────────────────────────────────

def _extract_pdf_text_fitz(data: bytes) -> str:
    """Extrae texto de un PDF usando pymupdf (fitz). Rápido y preciso."""
    try:
        import fitz  # pymupdf
        doc = fitz.open(stream=data, filetype="pdf")
        parts: list[str] = []
        for page in doc:
            text = page.get_text("text")
            if text and text.strip():
                parts.append(text.strip())
        doc.close()
        return "\n\n".join(parts)
    except Exception as e:
        logger.debug("      fitz error: %s", e)
        return ""


def _rasterize_pdf_pages(data: bytes, max_pages: int = RASTER_MAX_PAGES, dpi: int = RASTER_DPI) -> list[bytes]:
    """
    Convierte las primeras páginas de un PDF escaneado en imágenes JPEG.
    Permite que Claude "vea" PDFs gigantes (clippings de prensa de 50MB)
    que no caben como bloque nativo.
    """
    try:
        import fitz  # pymupdf
        doc = fitz.open(stream=data, filetype="pdf")
        out: list[bytes] = []
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            pix = page.get_pixmap(dpi=dpi)
            try:
                out.append(pix.tobytes("jpeg"))
            except Exception:
                out.append(pix.tobytes("png"))
        doc.close()
        return out
    except Exception as e:
        logger.debug("      rasterize error: %s", e)
        return []


def _extract_pdf_text_pdfminer(data: bytes) -> str:
    """Extrae texto de un PDF usando pdfminer.six como fallback de fitz."""
    try:
        from pdfminer.high_level import extract_text as pm_extract
        text = pm_extract(io.BytesIO(data))
        return (text or "").strip()
    except Exception as e:
        logger.debug("      pdfminer error: %s", e)
        return ""


# ─────────────────────────────────────────────────────────────────────────────
# Extractores de texto para formatos Office
# ─────────────────────────────────────────────────────────────────────────────

def _extract_docx(data: bytes) -> str:
    import docx  # python-docx

    doc = docx.Document(io.BytesIO(data))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation  # python-pptx

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        notes = getattr(slide, "notes_slide", None)
        if notes and notes.notes_text_frame and notes.notes_text_frame.text.strip():
            parts.append("Notas: " + notes.notes_text_frame.text.strip())
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    out = io.StringIO()
    writer = csv.writer(out)
    for ws in wb.worksheets:
        out.write(f"### Hoja: {ws.title}\n")
        rows_written = 0
        for row in ws.iter_rows(values_only=True):
            if rows_written >= 500:  # tope por hoja para no explotar
                out.write("…[más filas omitidas]…\n")
                break
            if any(c is not None for c in row):
                writer.writerow(["" if c is None else c for c in row])
                rows_written += 1
        out.write("\n")
    wb.close()
    return out.getvalue()


def _extract_email_like(data: bytes) -> str:
    """
    Extrae texto legible de un correo .eml o un .mhtml (página web guardada,
    p.ej. transcripciones exportadas). Ambos son mensajes MIME.
    """
    import email
    import email.policy

    msg = email.message_from_bytes(data, policy=email.policy.default)
    parts: list[str] = []
    for header in ("From", "To", "Cc", "Date", "Subject"):
        if msg.get(header):
            parts.append(f"{header}: {msg[header]}")
    if parts:
        parts.append("")

    plain: list[str] = []
    html: list[str] = []
    for part in msg.walk():
        ctype = part.get_content_type()
        if ctype not in ("text/plain", "text/html"):
            continue
        try:
            content = part.get_content()
        except Exception:
            payload = part.get_payload(decode=True)
            content = payload.decode("utf-8", errors="replace") if payload else ""
        if not content:
            continue
        if ctype == "text/plain":
            plain.append(content)
        else:
            html.append(_strip_html(content))

    # Preferir texto plano; el HTML suele duplicar lo mismo con markup
    parts.extend(plain if plain else html)
    return "\n".join(parts).strip()


def _strip_html(html: str) -> str:
    import html as html_mod
    import re

    text = re.sub(r"<(script|style)[^>]*>.*?</\1>", " ", html, flags=re.S | re.I)
    text = re.sub(r"<(br|/p|/div|/tr|/li|/h[1-6])[^>]*>", "\n", text, flags=re.I)
    text = re.sub(r"<[^>]+>", " ", text)
    text = html_mod.unescape(text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n\s*\n+", "\n\n", text)
    return text.strip()


def _extract_zip_text(data: bytes) -> str:
    """Lee los .txt dentro de un .zip (típicamente la exportación de WhatsApp).
    Alias de _extract_zip_text_raw — usado por _file_to_block para ZIPs fuera de la
    subcarpeta 04 o cuando wa_context no está disponible."""
    return _extract_zip_text_raw(data)


def _extract_zip_text_raw(data: bytes) -> str:
    """Extrae TODOS los .txt de un ZIP sin filtrar (lectura completa, sin watermark)."""
    parts: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            txt_names = [n for n in zf.namelist() if n.lower().endswith(".txt")]
            for n in txt_names:
                with zf.open(n) as fh:
                    parts.append(f"--- {n} ---")
                    parts.append(fh.read().decode("utf-8", errors="replace"))
            if not txt_names:
                return "(zip sin archivos de texto; contiene: " + ", ".join(zf.namelist()[:10]) + ")"
    except zipfile.BadZipFile:
        return "(archivo zip ilegible)"
    return "\n".join(parts)


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def _to_int(v) -> int | None:
    try:
        return int(v)
    except (TypeError, ValueError):
        return None


def _mb(n: int) -> str:
    return f"{n / 1024 / 1024:.1f} MB"
